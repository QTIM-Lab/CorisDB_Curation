import pdb
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
from pathlib import Path

from DICOMParser import DICOMParser                                                                                                                                                                                                                                                                                                                     

import pandas as pd

def create_image_presentation(csv_df, output_pptx):
    """
    Creates a PowerPoint presentation with image previews
    """
    # Create presentation object
    prs = Presentation()
    # Get slide dimensions (default is 10" x 7.5")
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    # Image dimensions (adjust as needed)
    img_width = Inches(3)
    img_height = Inches(3)
    # Get all image files
    image_extensions = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif'}
    # csv_df_test = csv_df['file_path'].tolist()[0:2]
    # [(row[1]['file_path'],row[1]['QTIM_Modality'],row[1]['Misc']) for row in csv_df[['file_path','QTIM_Modality','Misc']].iloc[0:2].iterrows()]
    # for img_file in sorted(image_files):
    tuples_list = csv_df[['Manufacturer', 'ManufacturerModelName', 'PhotometricInterpretation', 'SeriesDescription', 'QTIM_Modality']].drop_duplicates().values.tolist()
    def match_or_na(series, value):
        """Return True where series == value OR both are NaN."""
        return (pd.isna(value) & series.isna()) | series.eq(value)

    for Manufacturer, ManufacturerModelName, PhotometricInterpretation, SeriesDescription, QTIM_Modality in tuples_list:
        three_rows = csv_df[
            match_or_na(csv_df['Manufacturer'], Manufacturer) &
            match_or_na(csv_df['ManufacturerModelName'], ManufacturerModelName) &
            match_or_na(csv_df['PhotometricInterpretation'], PhotometricInterpretation) &
            match_or_na(csv_df['SeriesDescription'], SeriesDescription) &
            match_or_na(csv_df['QTIM_Modality'], QTIM_Modality)
        ]
        # if three_rows.shape[0] != 3:
        #     pdb.set_trace()
        try:
            one = three_rows.iloc[0]
        except:
            print("First image not working")
        try:
            two = three_rows.iloc[1]
        except:
            print("No second image, not failing just yet")
            two = None
        try:
            three = three_rows.iloc[2]
        except:
            print("No third image, not failing just yet")
            three = None
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        # 1: The horizontal position (left margin) of the title box from the left edge of the slide.
        # 2: The vertical position (top margin) of the title box from the top edge of the slide.
        # 3: The width of the title box.
        # 4: The height of the title box.
        top_title = Inches(0.2)
        top_image = Inches(1.2)
        top_subtitle = Inches(0.5)
        top_subtitle_below_image = Inches(5)
        # pdb.set_trace()
        ## Add title 1
        col_1 = slide_width - img_width*3 - Inches(0.1*3) - Inches(0.3)
        title_box_1 = slide.shapes.add_textbox(col_1,top_title,Inches(2),Inches(0.8))
        title_frame_1 = title_box_1.text_frame
        title_frame_1.text = f"{one['QTIM_Modality']} ({one['count']})"
        title_frame_1.paragraphs[0].font.size = Pt(10)
        title_frame_1.paragraphs[0].alignment = PP_ALIGN.LEFT
        ## Add subtitle 1
        subtitle_box_1 = slide.shapes.add_textbox(col_1,top_subtitle,Inches(2),Inches(0.8))
        subtitle_frame_1 = subtitle_box_1.text_frame
        subtitle_frame_1.text = f"Misc: {one['Misc']}"
        subtitle_frame_1.paragraphs[0].font.size = Pt(10)
        subtitle_frame_1.paragraphs[0].alignment = PP_ALIGN.LEFT
        ## Add subtitle 1 - below image
        subtitle_box_1_below_image = slide.shapes.add_textbox(col_1,top_subtitle_below_image,Inches(2),Inches(0.8))
        subtitle_frame_1_below_image = subtitle_box_1_below_image.text_frame
        subtitle_frame_1_below_image.text = f"""
        file_path: {one['file_path']}
        SOPClassDescription: {one['SOPClassDescription']}
        Modality: {one['Modality']}
        Manufacturer: {one['Manufacturer']}
        ManufacturerModelName: {one['ManufacturerModelName']}
        PhotometricInterpretation: {one['PhotometricInterpretation']}
        SeriesDescription: {one['SeriesDescription']}
        """
        for p in subtitle_frame_1_below_image.paragraphs:
            for run in p.runs:
                run.font.size = Pt(8)
        for p in subtitle_frame_1_below_image.paragraphs:
            for run in p.runs:
                run.slignment = PP_ALIGN.LEFT
        # Add image (centered)
        left_1 = col_1
        top_1 = top_image
        ## Add title 2
        if two is not None:
            col_2 = slide_width - img_width*2 - Inches(0.1*2) - Inches(0.3)
            title_box_2 = slide.shapes.add_textbox(col_2,top_title,Inches(2),Inches(0.8))
            title_frame_2 = title_box_2.text_frame
            title_frame_2.text = f"{two['QTIM_Modality']} ({two['count']})"
            title_frame_2.paragraphs[0].font.size = Pt(10)
            title_frame_2.paragraphs[0].alignment = PP_ALIGN.LEFT
            ## Add subtitle 2
            subtitle_box_2 = slide.shapes.add_textbox(col_2,top_subtitle + Inches(0.1),Inches(2),Inches(0.8))
            subtitle_frame_2 = subtitle_box_2.text_frame
            subtitle_frame_2.text = f"Misc: {two['Misc']}"
            subtitle_frame_2.paragraphs[0].font.size = Pt(10)
            subtitle_frame_2.paragraphs[0].alignment = PP_ALIGN.LEFT
            ## Add subtitle 2 - below image
            subtitle_box_2_below_image = slide.shapes.add_textbox(col_2,top_subtitle_below_image,Inches(2),Inches(0.8))
            subtitle_frame_2_below_image = subtitle_box_2_below_image.text_frame
            subtitle_frame_2_below_image.text = f"""


            

            file_path: {two['file_path']}
            SOPClassDescription: {two['SOPClassDescription']}
            Modality: {two['Modality']}
            Manufacturer: {two['Manufacturer']}
            ManufacturerModelName: {two['ManufacturerModelName']}
            PhotometricInterpretation: {two['PhotometricInterpretation']}
            SeriesDescription: {two['SeriesDescription']}
            """
            for p in subtitle_frame_2_below_image.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(8)
            for p in subtitle_frame_2_below_image.paragraphs:
                for run in p.runs:
                    run.slignment = PP_ALIGN.LEFT
            # Add image (centered)
            left_2 = col_2
            top_2 = top_image
        ## Add title 3
        if three is not None:
            col_3 = slide_width - img_width*1 - Inches(0.1*1) - Inches(0.3)
            title_box_3 = slide.shapes.add_textbox(col_3,top_title,Inches(2),Inches(0.8))
            title_frame_3 = title_box_3.text_frame
            title_frame_3.text = f"{three['QTIM_Modality']} ({three['count']})"
            title_frame_3.paragraphs[0].font.size = Pt(10)
            title_frame_3.paragraphs[0].alignment = PP_ALIGN.LEFT
            ## Add subtitle 3
            subtitle_box_3 = slide.shapes.add_textbox(col_3,top_subtitle + Inches(0.2),Inches(2),Inches(0.8))
            subtitle_frame_3 = subtitle_box_3.text_frame
            subtitle_frame_3.text = f"Misc: {three['Misc']}"
            subtitle_frame_3.paragraphs[0].font.size = Pt(10)
            subtitle_frame_3.paragraphs[0].alignment = PP_ALIGN.CENTER
            ## Add subtitle 3 - below image
            subtitle_box_3_below_image = slide.shapes.add_textbox(col_3,top_subtitle_below_image,Inches(2),Inches(0.8))
            subtitle_frame_3_below_image = subtitle_box_3_below_image.text_frame
            subtitle_frame_3_below_image.text = f"""
            file_path: {three['file_path']}
            SOPClassDescription: {three['SOPClassDescription']}
            Modality: {three['Modality']}
            Manufacturer: {three['Manufacturer']}
            ManufacturerModelName: {three['ManufacturerModelName']}
            PhotometricInterpretation: {three['PhotometricInterpretation']}
            SeriesDescription: {three['SeriesDescription']}
            """
            for p in subtitle_frame_3_below_image.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(8)
            for p in subtitle_frame_3_below_image.paragraphs:
                for run in p.runs:
                    run.slignment = PP_ALIGN.LEFT
            # Add image (centered)
            left_3 = col_3
            top_3 = top_image
        try:
            parser_1 = DICOMParser.create_parser(one['file_path'])
            parsed_1 = parser_1.parse()
            image_stream_1 = BytesIO()
            if two is not None:
                parser_2 = DICOMParser.create_parser(two['file_path'])
                image_stream_2 = BytesIO()
                parsed_2 = parser_2.parse()
            if three is not None:
                parser_3 = DICOMParser.create_parser(three['file_path'])
                parsed_3 = parser_3.parse()
                image_stream_3 = BytesIO()
            try:
                parsed_1['image_PIL'].save(image_stream_1, format='PNG')  # Save the PIL image to a BytesIO stream
                image_stream_1.seek(0)  # Reset stream position to the beginning
                slide.shapes.add_picture(image_stream_1, left_1, top_1, img_width, img_height)
                if two is not None:
                    parsed_2['image_PIL'].save(image_stream_2, format='PNG')  # Save the PIL image to a BytesIO stream
                    image_stream_2.seek(0)  # Reset stream position to the beginning
                    slide.shapes.add_picture(image_stream_2, left_2, top_2, img_width, img_height)
                if three is not None:
                    parsed_3['image_PIL'].save(image_stream_3, format='PNG')  # Save the PIL image to a BytesIO stream
                    image_stream_3.seek(0)  # Reset stream position to the beginning
                    slide.shapes.add_picture(image_stream_3, left_3, top_3, img_width, img_height)
            except:
                pdb.set_trace()
                try:
                    if 'png_pages' in parsed_1:
                        parsed_1['png_pages']['page_PIL']['page_1'].save(image_stream_1, format='PNG')  # Save the PIL image to a BytesIO stream
                    else:
                        parsed_1['bscan_images']['bscan1'].save(image_stream_1, format='PNG')  # Save the PIL image to a BytesIO stream
                    image_stream_1.seek(0)  # Reset stream position to the beginning
                    slide.shapes.add_picture(image_stream_1, left_1, top_1, img_width, img_height)
                    pdb.set_trace()
                    if two is not None:
                        if 'png_pages' in parsed_2:
                            parsed_2['png_pages']['page_PIL']['page_1'].save(image_stream_2, format='PNG')  # Save the PIL image to a BytesIO stream
                        else:
                            parsed_2['bscan_images']['bscan1'].save(image_stream_2, format='PNG')  # Save the PIL image to a BytesIO stream
                        image_stream_2.seek(0)  # Reset stream position to the beginning
                        slide.shapes.add_picture(image_stream_2, left_2, top_2, img_width, img_height)
                    pdb.set_trace()
                    if three is not None:
                        if 'png_pages' in parsed_3:
                            parsed_3['png_pages']['page_PIL']['page_1']['page_PIL'].save(image_stream_3, format='PNG')  # Save the PIL image to a BytesIO stream
                        else:
                            parsed_3['bscan_images']['bscan1']['page_PIL'].save(image_stream_3, format='PNG')  # Save the PIL image to a BytesIO stream
                        image_stream_3.seek(0)  # Reset stream position to the beginning
                        slide.shapes.add_picture(image_stream_3, left_3, top_3, img_width, img_height)
                    pdb.set_trace()
                except:
                    # Format the error text
                    # --- Error box 1 ---
                    error_text_1 = "What was parsed by DICOMParser.py:\n" + "\n".join([f"{key}: {value}" for key, value in parsed_1.items()])
                    error_box_1 = slide.shapes.add_textbox(left_1, top_1, img_width, img_height)
                    error_frame_1 = error_box_1.text_frame
                    error_frame_1.text = error_text_1
                    for p in error_frame_1.paragraphs:
                        p.alignment = PP_ALIGN.LEFT
                        for run in p.runs:
                            run.font.size = Pt(8)  # or whatever size makes sense for error text
                    if two is not None:
                        error_text_2 = "What was parsed by DICOMParser.py:\n" + "\n".join([f"{key}: {value}" for key, value in parsed_2.items()])
                        # --- Error box 2 ---
                        error_box_2 = slide.shapes.add_textbox(left_2, top_2+Inches(1.3), img_width, img_height)
                        error_frame_2 = error_box_2.text_frame
                        error_frame_2.text = error_text_2
                        for p in error_frame_2.paragraphs:
                            p.alignment = PP_ALIGN.LEFT
                            for run in p.runs:
                                run.font.size = Pt(8)
                    if three is not None:
                        error_text_3 = "What was parsed by DICOMParser.py:\n" + "\n".join([f"{key}: {value}" for key, value in parsed_3.items()])
                        # --- Error box 3 ---
                        error_box_3 = slide.shapes.add_textbox(left_3, top_3+Inches(2.6), img_width, img_height)
                        error_frame_3 = error_box_3.text_frame
                        error_frame_3.text = error_text_3
                        for p in error_frame_3.paragraphs:
                            p.alignment = PP_ALIGN.LEFT
                            for run in p.runs:
                                run.font.size = Pt(8)
        except Exception as e:
            pdb.set_trace()
            print(f"Error adding {one['file_path']}: {e} or\n")
            print(f"Error adding {two['file_path']}: {e} or\n")
            print(f"Error adding {three['file_path']}: {e}")
            # Add error text instead
            error_box_1 = slide.shapes.add_textbox(left_1, top_1, img_width, img_height)
            error_frame_1 = error_box_1.text_frame
            error_frame_1.text = f"Could not load image:\n{one['file_path']}"
            error_frame_1.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Add error text instead
            error_box_2 = slide.shapes.add_textbox(left_2, top_2, img_width, img_height)
            error_frame_2 = error_box_2.text_frame
            error_frame_2.text = f"Could not load image:\n{two['file_path']}"
            error_frame_2.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Add error text instead
            error_box_3 = slide.shapes.add_textbox(left_3, top_3, img_width, img_height)
            error_frame_3 = error_box_3.text_frame
            error_frame_3.text = f"Could not load image:\n{three['file_path']}"
            error_frame_3.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Save presentation
    prs.save(os.path.join(output_pptx))
    print(f"Presentation saved as {os.path.join(output_pptx)}")



# Usage example
if __name__ == "__main__":
    csv = "/scratch90/QTIM/Active/23-0284/EHR/AXISPACS/col_counts/dicoms_preview.csv"
    output_pptx = "/scratch90/QTIM/Active/23-0284/EHR/AXISPACS/col_counts/dicoms_preview.pptx"
    csv_df = pd.read_csv(csv)
    create_image_presentation(csv_df, output_pptx)