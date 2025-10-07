import os, pdb
import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

from preview import preview


def create_image_presentation(csv_df, output_pptx):
    """
    Creates a PowerPoint presentation with image previews
    """
    # Create presentation object
    prs = Presentation()
    # Get slide dimensions (default is 10" x 7.5")
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    # Image dimensions (adjust as needed)
    img_width = Inches(3)
    img_height = Inches(3)
    # Get all image files
    image_extensions = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif'}
    # csv_df_test = csv_df['file_path'].tolist()[0:2]
    # [(row[1]['file_path'],row[1]['QTIM_Modality'],row[1]['Misc']) for row in csv_df[['file_path','QTIM_Modality','Misc']].iloc[0:2].iterrows()]
    # for img_file in sorted(image_files):
    tuples_list = csv_df[['ParsedImageGroup', 'ReportType', 'Layer_Name', 'Procedure', 'ScanPattern', 'QTIM_Modality']].drop_duplicates().values.tolist()
    for ParsedImageGroup, ReportType, Layer_Name, Procedure, ScanPattern, QTIM_Modality in tuples_list:
        three_rows = csv_df[
                       ((csv_df['ParsedImageGroup'] == ParsedImageGroup) | (pd.isna(csv_df['ParsedImageGroup']))) &
                       ((csv_df['ReportType'] == ReportType) | (pd.isna(csv_df['ReportType']))) &
                       ((csv_df['Layer_Name'] == Layer_Name) | (pd.isna(csv_df['Layer_Name']))) &
                       ((csv_df['Procedure'] == Procedure) | (pd.isna(csv_df['Procedure']))) &
                       ((csv_df['ScanPattern'] == ScanPattern) | (pd.isna(csv_df['ScanPattern']))) &
                       ((csv_df['QTIM_Modality'] == QTIM_Modality) | (pd.isna(csv_df['QTIM_Modality'])))
                      ]
        try:
            one = three_rows.iloc[0]
        except:
            print("First image not working")
            pdb.set_trace()
        try:
            two = three_rows.iloc[1]
        except:
            print("No second image, not failing just yet")
            two = None
        try:
            three = three_rows.iloc[2]
            print("No second image, not failing just yet")
        except:
            print("No third image, not failing just yet")
            three = None
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        ## Add title
        # 1: The horizontal position (left margin) of the title box from the left edge of the slide.
        # 2: The vertical position (top margin) of the title box from the top edge of the slide.
        # 3: The width of the title box.
        # 4: The height of the title box.
        top_title = Inches(0.2)
        top_image = Inches(1.2)
        top_subtitle = Inches(0.5)
        top_subtitle_below_image = Inches(5)
        # pdb.set_trace()
        ## Add title 1
        col_1 = slide_width - img_width*3 - Inches(0.1*3) - Inches(0.3)
        title_box_1 = slide.shapes.add_textbox(col_1,top_title,Inches(2),Inches(0.8))
        title_frame_1 = title_box_1.text_frame
        title_frame_1.text = f"{one['QTIM_Modality']} ({one['count']})"
        title_frame_1.paragraphs[0].font.size = Pt(8)
        title_frame_1.paragraphs[0].alignment = PP_ALIGN.LEFT
        ## Add subtitle 1
        subtitle_box_1 = slide.shapes.add_textbox(col_1,top_subtitle,Inches(2),Inches(0.8))
        subtitle_frame_1 = subtitle_box_1.text_frame
        subtitle_frame_1.text = f"Misc: {one['Misc']}"
        subtitle_frame_1.paragraphs[0].font.size = Pt(8)
        subtitle_frame_1.paragraphs[0].alignment = PP_ALIGN.LEFT
        ## Add subtitle 1 - below image
        subtitle_box_1_below_image = slide.shapes.add_textbox(col_1,top_subtitle_below_image,Inches(2),Inches(0.8))
        subtitle_frame_1_below_image = subtitle_box_1_below_image.text_frame
        subtitle_frame_1_below_image.text = f"""
        file_path: {one['file_path']}
        ParsedImageGroup: {one['ParsedImageGroup']}
        ReportType: {one['ReportType']}
        Layer_Name: {one['Layer_Name']}
        Procedure: {one['Procedure']}
        ScanPattern: {one['ScanPattern']}
        QTIM_Modality: {one['QTIM_Modality']}
        """
        for p in subtitle_frame_1_below_image.paragraphs:
            for run in p.runs:
                run.font.size = Pt(8)
        for p in subtitle_frame_1_below_image.paragraphs:
            for run in p.runs:
                run.slignment = PP_ALIGN.LEFT
        # Add image (centered)
        left_1 = col_1
        top_1 = top_image
        ## Add title 2
        if two is not None:
            col_2 = slide_width - img_width*2 - Inches(0.1*2) - Inches(0.3)
            title_box_2 = slide.shapes.add_textbox(col_2,top_title,Inches(2),Inches(0.8))
            title_frame_2 = title_box_2.text_frame
            title_frame_2.text = f"{two['QTIM_Modality']} ({two['count']})"
            title_frame_2.paragraphs[0].font.size = Pt(8)
            title_frame_2.paragraphs[0].alignment = PP_ALIGN.LEFT
            ## Add subtitle 2
            subtitle_box_2 = slide.shapes.add_textbox(col_2,top_subtitle + Inches(0.1),Inches(2),Inches(0.8))
            subtitle_frame_2 = subtitle_box_2.text_frame
            subtitle_frame_2.text = f"Misc: {two['Misc']}"
            subtitle_frame_2.paragraphs[0].font.size = Pt(8)
            subtitle_frame_2.paragraphs[0].alignment = PP_ALIGN.LEFT
            ## Add subtitle 2 - below image
            subtitle_box_2_below_image = slide.shapes.add_textbox(col_2,top_subtitle_below_image,Inches(2),Inches(0.8))
            subtitle_frame_2_below_image = subtitle_box_2_below_image.text_frame
            subtitle_frame_2_below_image.text = f"""



            
            file_path: {two['file_path']}
            ParsedImageGroup: {two['ParsedImageGroup']}
            ReportType: {two['ReportType']}
            Layer_Name: {two['Layer_Name']}
            Procedure: {two['Procedure']}
            ScanPattern: {two['ScanPattern']}
            QTIM_Modality: {two['QTIM_Modality']}
            """
            for p in subtitle_frame_2_below_image.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(8)
            for p in subtitle_frame_2_below_image.paragraphs:
                for run in p.runs:
                    run.slignment = PP_ALIGN.LEFT
            # Add image (centered)
            left_2 = col_2
            top_2 = top_image
        ## Add title 3
        if three is not None:
            col_3 = slide_width - img_width*1 - Inches(0.1*1) - Inches(0.3)
            title_box_3 = slide.shapes.add_textbox(col_3,top_title,Inches(2),Inches(0.8))
            title_frame_3 = title_box_3.text_frame
            title_frame_3.text = f"{three['QTIM_Modality']} ({three['count']})"
            title_frame_3.paragraphs[0].font.size = Pt(8)
            title_frame_3.paragraphs[0].alignment = PP_ALIGN.LEFT
            ## Add subtitle 3
            subtitle_box_3 = slide.shapes.add_textbox(col_3,top_subtitle + Inches(0.2),Inches(2),Inches(0.8))
            subtitle_frame_3 = subtitle_box_3.text_frame
            subtitle_frame_3.text = f"Misc: {three['Misc']}"
            subtitle_frame_3.paragraphs[0].font.size = Pt(8)
            subtitle_frame_3.paragraphs[0].alignment = PP_ALIGN.CENTER
            ## Add subtitle 3 - below image
            subtitle_box_3_below_image = slide.shapes.add_textbox(col_3,top_subtitle_below_image,Inches(2),Inches(0.8))
            subtitle_frame_3_below_image = subtitle_box_3_below_image.text_frame
            subtitle_frame_3_below_image.text = f"""
            file_path: {three['file_path']}
            ParsedImageGroup: {three['ParsedImageGroup']}
            ReportType: {three['ReportType']}
            Layer_Name: {three['Layer_Name']}
            Procedure: {three['Procedure']}
            ScanPattern: {three['ScanPattern']}
            QTIM_Modality: {three['QTIM_Modality']}
            """
            for p in subtitle_frame_3_below_image.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(8)
            for p in subtitle_frame_3_below_image.paragraphs:
                for run in p.runs:
                    run.slignment = PP_ALIGN.LEFT
            # Add image (centered)
            left_3 = col_3
            top_3 = top_image
        try:
            parsed_1 = preview(one['file_path'], None, return_img=True)
            image_stream_1 = BytesIO()
            if two is not None:
                parsed_2 = preview(two['file_path'], None, return_img=True)
                image_stream_2 = BytesIO()
            if three is not None:
                parsed_3 = preview(three['file_path'], None, return_img=True)
                image_stream_3 = BytesIO()
            try:
                parsed_1.save(image_stream_1, format='PNG')  # Save the PIL image to a BytesIO stream
                image_stream_1.seek(0)  # Reset stream position to the beginning
                slide.shapes.add_picture(image_stream_1, left_1, top_1, img_width, img_height)
                if two is not None:
                    parsed_2.save(image_stream_2, format='PNG')  # Save the PIL image to a BytesIO stream
                    image_stream_2.seek(0)  # Reset stream position to the beginning
                    slide.shapes.add_picture(image_stream_2, left_2, top_2, img_width, img_height)
                if three is not None:
                    parsed_3.save(image_stream_3, format='PNG')  # Save the PIL image to a BytesIO stream
                    image_stream_3.seek(0)  # Reset stream position to the beginning
                    slide.shapes.add_picture(image_stream_3, left_3, top_3, img_width, img_height)
            except:
                pdb.set_trace()
        except Exception as e:
            pdb.set_trace()
            print(f"Error adding {one['file_path']}: {e} or\n")
            print(f"Error adding {two['file_path']}: {e} or\n")
            print(f"Error adding {three['file_path']}: {e}")
            error_box_1 = slide.shapes.add_textbox(left_1, top_1, img_width, img_height)
            error_frame_1 = error_box_1.text_frame
            error_frame_1.text = f"Could not load image:\n{one['file_path']}"
            error_frame_1.paragraphs[0].alignment = PP_ALIGN.CENTER
            if two is not None:
                error_box_2 = slide.shapes.add_textbox(left_2, top_2, img_width, img_height)
                error_frame_2 = error_box_2.text_frame
                error_frame_2.text = f"Could not load image:\n{two['file_path']}"
                error_frame_2.paragraphs[0].alignment = PP_ALIGN.CENTER
            if three is not None:
                error_box_3 = slide.shapes.add_textbox(left_3, top_3, img_width, img_height)
                error_frame_3 = error_box_3.text_frame
                error_frame_3.text = f"Could not load image:\n{three['file_path']}"
                error_frame_3.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Save presentation
    prs.save(os.path.join(output_pptx))
    print(f"Presentation saved as {os.path.join(output_pptx)}")



# Usage example
if __name__ == "__main__":
    csv = "/scratch90/QTIM/Active/23-0284/EHR/AXISPACS/col_counts/xmls_preview.csv"
    output_pptx = "/scratch90/QTIM/Active/23-0284/EHR/AXISPACS/col_counts/j2ks_preview.pptx"
    csv_df = pd.read_csv(csv)
    create_image_presentation(csv_df, output_pptx)